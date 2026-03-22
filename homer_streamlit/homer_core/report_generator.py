# Homer - Report Generator (PDF + PPTX)
# Creates downloadable reports with dashboard figures in grid layouts

import io
import math
import datetime
from typing import Optional

import plotly.graph_objects as go
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
import pandas as pd

from . import plotting


def _fig_to_png_bytes(plotly_fig: go.Figure, width: int = 1100,
                      height: int = 650, scale: int = 2) -> bytes:
    """Render a Plotly figure to PNG bytes."""
    return plotly_fig.to_image(format="png", width=width, height=height, scale=scale)


def _fig_to_compressed_png(plotly_fig: go.Figure, width: int = 1100,
                           height: int = 650, quality: int = 85) -> bytes:
    """Render Plotly figure to compressed PNG via Pillow."""
    from PIL import Image
    raw = plotly_fig.to_image(format="png", width=width, height=height, scale=2)
    img = Image.open(io.BytesIO(raw))
    if img.mode == "RGBA":
        bg = Image.new("RGB", img.size, (255, 255, 255))
        bg.paste(img, mask=img.split()[3])
        img = bg
    buf = io.BytesIO()
    img.save(buf, format="PNG", optimize=True)
    buf.seek(0)
    return buf.read()


class ReportBuilder:
    """Builds multi-page PDF or PPTX reports from dashboard figures."""

    def __init__(self, title: str = "Homer Data Report", dataset_name: str = ""):
        self.title = title
        self.dataset_name = dataset_name
        self.figures: list[dict] = []

    def add_figure(self, title: str, plotly_fig: go.Figure,
                   config: Optional[dict] = None):
        self.figures.append({
            "title": title,
            "plotly_fig": plotly_fig,
            "config": config or {},
        })

    def _create_title_page(self) -> plt.Figure:
        fig, ax = plt.subplots(figsize=(11, 8.5))
        ax.axis("off")
        ax.text(0.5, 0.7, "HOMER", transform=ax.transAxes,
                fontsize=48, fontweight="bold", ha="center", va="center",
                color="#2E86AB")
        ax.text(0.5, 0.60, "Histology Output Mapper & Explorer for Research",
                transform=ax.transAxes, fontsize=16, ha="center", va="center",
                color="#555555")
        ax.text(0.5, 0.45, self.title, transform=ax.transAxes,
                fontsize=20, ha="center", va="center", color="#333333")
        if self.dataset_name:
            ax.text(0.5, 0.38, f"Dataset: {self.dataset_name}",
                    transform=ax.transAxes, fontsize=14, ha="center", va="center",
                    color="#666666")
        ax.text(0.5, 0.25, f"Generated: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}",
                transform=ax.transAxes, fontsize=12, ha="center", va="center",
                color="#888888")
        ax.text(0.5, 0.20, f"Total figures: {len(self.figures)}",
                transform=ax.transAxes, fontsize=12, ha="center", va="center",
                color="#888888")
        fig.patch.set_facecolor("white")
        return fig

    def _plotly_to_mpl_image(self, plotly_fig: go.Figure) -> plt.Figure:
        try:
            img_bytes = plotly_fig.to_image(format="png", width=1100, height=650, scale=2)
            from PIL import Image
            img = Image.open(io.BytesIO(img_bytes))
            fig, ax = plt.subplots(figsize=(11, 7))
            ax.imshow(img)
            ax.axis("off")
            fig.patch.set_facecolor("white")
            fig.tight_layout(pad=0.5)
            return fig
        except Exception:
            fig, ax = plt.subplots(figsize=(11, 7))
            ax.text(0.5, 0.5, "Figure could not be rendered\n(install kaleido for image export)",
                    transform=ax.transAxes, fontsize=14, ha="center", va="center")
            ax.axis("off")
            return fig

    def generate_pdf(self) -> bytes:
        buf = io.BytesIO()
        with PdfPages(buf) as pdf:
            title_fig = self._create_title_page()
            pdf.savefig(title_fig, dpi=150)
            plt.close(title_fig)
            for entry in self.figures:
                fig = self._plotly_to_mpl_image(entry["plotly_fig"])
                fig.suptitle(entry["title"], fontsize=14, fontweight="bold", y=0.98)
                pdf.savefig(fig, dpi=150)
                plt.close(fig)
        buf.seek(0)
        return buf.read()

    def generate_pptx(self, grid_cols: int = 2, include_data_table: bool = False,
                      df: Optional[pd.DataFrame] = None) -> bytes:
        """Generate PPTX with figures arranged in a grid layout.

        Args:
            grid_cols: Number of columns in the figure grid (1-4).
            include_data_table: If True, append a slide with a data table.
            df: DataFrame for the data table slide.
        """
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        # ── Title slide ──
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "HOMER"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x4F, 0xC3, 0xF7)
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = "Histology Output Mapper & Explorer for Research"
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
        p2.alignment = PP_ALIGN.CENTER

        p3 = tf.add_paragraph()
        p3.text = self.title
        p3.font.size = Pt(24)
        p3.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(24)

        p4 = tf.add_paragraph()
        p4.text = f"{self.dataset_name}  |  {datetime.datetime.now().strftime('%Y-%m-%d')}  |  {len(self.figures)} figures"
        p4.font.size = Pt(14)
        p4.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
        p4.alignment = PP_ALIGN.CENTER
        p4.space_before = Pt(12)

        # ── Figure slides in grid ──
        margin = Inches(0.4)
        title_h = Inches(0.6)
        grid_cols = max(1, min(4, grid_cols))
        grid_rows = 2 if grid_cols > 1 else 1
        figs_per_slide = grid_cols * grid_rows

        cell_w = (slide_w - margin * 2) // grid_cols
        cell_h = (slide_h - margin - title_h) // grid_rows
        img_w = int(cell_w - Inches(0.2))
        img_h = int(cell_h - Inches(0.5))

        # Render all figures to PNG
        rendered = []
        for entry in self.figures:
            try:
                png = _fig_to_compressed_png(
                    entry["plotly_fig"],
                    width=max(400, int(img_w / 914400 * 96)),
                    height=max(300, int(img_h / 914400 * 96)),
                )
                rendered.append({"title": entry["title"], "png": png})
            except Exception:
                rendered.append({"title": entry["title"], "png": None})

        n_slides = math.ceil(len(rendered) / figs_per_slide) if rendered else 0
        for slide_idx in range(n_slides):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            batch = rendered[slide_idx * figs_per_slide:(slide_idx + 1) * figs_per_slide]
            for i, item in enumerate(batch):
                col = i % grid_cols
                row = i // grid_cols
                x = int(margin) + col * int(cell_w)
                y = int(title_h) + row * int(cell_h)

                # Title
                tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(int(cell_w)), Inches(0.4))
                p = tb.text_frame.paragraphs[0]
                p.text = item["title"]
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                p.alignment = PP_ALIGN.CENTER

                # Image
                if item["png"]:
                    img_stream = io.BytesIO(item["png"])
                    slide.shapes.add_picture(
                        img_stream,
                        Emu(x + int(Inches(0.1))),
                        Emu(y + int(Inches(0.4))),
                        Emu(img_w),
                        Emu(img_h),
                    )

        # ── Data table slide ──
        if include_data_table and df is not None and len(df) > 0:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            tb = slide.shapes.add_textbox(margin, Inches(0.2), Inches(12), Inches(0.5))
            p = tb.text_frame.paragraphs[0]
            p.text = "Data Summary"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

            # Limit to manageable size
            show_df = df.head(20)
            show_cols = list(show_df.columns[:10])
            n_rows = len(show_df) + 1
            n_cols = len(show_cols)

            tbl = slide.shapes.add_table(
                n_rows, n_cols,
                margin, Inches(0.9),
                Inches(min(12.5, n_cols * 1.5)), Inches(min(5.5, n_rows * 0.3)),
            ).table

            for j, col in enumerate(show_cols):
                cell = tbl.cell(0, j)
                cell.text = str(col)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(8)
                    paragraph.font.bold = True

            for i, (_, row) in enumerate(show_df[show_cols].iterrows()):
                for j, col in enumerate(show_cols):
                    cell = tbl.cell(i + 1, j)
                    val = row[col]
                    cell.text = f"{val:.2f}" if isinstance(val, float) else str(val)
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(7)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

    def clear(self):
        self.figures.clear()


def generate_data_summary_page(df: pd.DataFrame, data_type: str) -> plt.Figure:
    """Create a summary statistics page for the dataset."""
    fig, ax = plt.subplots(figsize=(11, 8.5))
    ax.axis("off")

    ax.text(0.5, 0.95, "Dataset Summary", transform=ax.transAxes,
            fontsize=20, fontweight="bold", ha="center", va="top", color="#2E86AB")

    summary_lines = [
        f"Data Type: {data_type.title()}",
        f"Rows: {len(df):,}",
        f"Columns: {len(df.columns)}",
        f"Numeric Columns: {len(df.select_dtypes(include='number').columns)}",
        f"Categorical Columns: {len(df.select_dtypes(include='object').columns)}",
        "",
    ]

    numeric_cols = df.select_dtypes(include="number").columns[:10]
    if len(numeric_cols) > 0:
        summary_lines.append("Numeric Column Statistics (first 10):")
        summary_lines.append(f"{'Column':<35} {'Mean':>12} {'Std':>12} {'Min':>12} {'Max':>12}")
        summary_lines.append("-" * 83)
        for col in numeric_cols:
            mean = f"{df[col].mean():.2f}" if pd.notna(df[col].mean()) else "N/A"
            std = f"{df[col].std():.2f}" if pd.notna(df[col].std()) else "N/A"
            mn = f"{df[col].min():.2f}" if pd.notna(df[col].min()) else "N/A"
            mx = f"{df[col].max():.2f}" if pd.notna(df[col].max()) else "N/A"
            col_display = col[:33] if len(col) > 33 else col
            summary_lines.append(f"{col_display:<35} {mean:>12} {std:>12} {mn:>12} {mx:>12}")

    text = "\n".join(summary_lines)
    ax.text(0.05, 0.85, text, transform=ax.transAxes, fontsize=9,
            fontfamily="monospace", va="top", color="#333333")

    fig.patch.set_facecolor("white")
    return fig
